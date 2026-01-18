"""
Text Extraction Module
Supports: PDF, DOCX, PPTX, Images (OCR)
"""

import os
import tempfile
import re
import html
from typing import Optional

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files"""
    try:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except ImportError:
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.strip()
        except ImportError:
            raise ImportError("Please install PyPDF2 or pdfplumber: pip install PyPDF2 pdfplumber")

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX files"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        raise ImportError("Please install python-docx: pip install python-docx")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except ImportError:
        raise ImportError("Please install python-pptx: pip install python-pptx")

def extract_text_from_image(file_path: str) -> str:
    """Extract text from images using OCR with multi-language support"""
    try:
        import pytesseract
        from PIL import Image
        
        # Configure tesseract path explicitly
        tesseract_path = r'D:\tesercat\tesseract.exe'
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        else:
            # Fallback to common Windows paths
            possible_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe'
            ]
            tesseract_found = False
            for path in possible_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    tesseract_found = True
                    break
            
            if not tesseract_found:
                raise Exception(f"Tesseract not found at {tesseract_path} or standard locations")
        
        # Open and validate image
        try:
            image = Image.open(file_path)
        except Exception as e:
            raise Exception(f"Failed to open image: {str(e)}")
        
        # Normalize image format for OCR
        if image.mode in ('RGBA', 'LA', 'P'):
            # Convert to RGB for better OCR results
            image = image.convert('RGB')
        elif image.mode == 'L':
            # Keep grayscale as is (good for OCR)
            pass
        elif image.mode != 'RGB':
            # Convert any other mode to RGB
            image = image.convert('RGB')
        
        # Try multiple language combinations for better text detection
        language_configs = [
            ('urd+eng', '--psm 6 -c tessedit_char_whitelist=اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789.,!?;:()"\'-'),  # Urdu + English with character whitelist
            ('urd', '--psm 6 -c tessedit_char_whitelist=اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے .,!?;:()"\'-'),      # Urdu only with whitelist
            ('urd+eng', '--psm 6'),  # Urdu + English fallback
            ('eng+urd', '--psm 6'),  # English + Urdu fallback
        ]
        
        best_text = ""
        best_confidence = 0
        
        for lang, config in language_configs:
            try:
                # Extract text with current language configuration
                text = pytesseract.image_to_string(image, lang=lang, config=config)
                
                if text and text.strip():
                    # Get confidence score
                    try:
                        data = pytesseract.image_to_data(image, lang=lang, config=config, output_type=pytesseract.Output.DICT)
                        confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
                        avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                        
                        if avg_confidence > best_confidence:
                            best_confidence = avg_confidence
                            best_text = text
                    except:
                        # If confidence calculation fails, use text length as a heuristic
                        if len(text.strip()) > len(best_text.strip()):
                            best_text = text
                            
            except Exception as e:
                # Continue with next language configuration
                continue
        
        if not best_text or not best_text.strip():
            raise Exception("No text detected in image with any language configuration")
        
        # Clean and normalize the extracted text
        cleaned_text = clean_extracted_text(best_text)
        
        return cleaned_text
    except ImportError:
        raise ImportError("Please install pytesseract and Pillow: pip install pytesseract Pillow")
    except Exception as e:
        raise Exception(f"OCR failed: {str(e)}")

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from DOC files (legacy Word format)"""
    try:
        import textract
        text = textract.process(file_path).decode('utf-8')
        return text.strip()
    except ImportError:
        raise ImportError("Please install textract: pip install textract")
    except Exception as e:
        # Fallback method using python-docx2txt
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            return text.strip() if text else ""
        except ImportError:
            raise ImportError("Please install docx2txt: pip install docx2txt")

def extract_text_from_file(file_path: str, file_extension: str) -> str:
    """
    Main function to extract text from various file types
    """
    file_extension = file_extension.lower()
    
    # Validate file path
    if not os.path.exists(file_path) or not os.path.isfile(file_path):
        raise ValueError("Invalid file path")
    
    # Check file size (max 50MB)
    if os.path.getsize(file_path) > 50 * 1024 * 1024:
        raise ValueError("File too large")
    
    try:
        if file_extension == '.txt':
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise Exception("Could not decode text file with any supported encoding")
        
        elif file_extension == '.pdf':
            return extract_text_from_pdf(file_path)
        
        elif file_extension == '.docx':
            return extract_text_from_docx(file_path)
        
        elif file_extension == '.doc':
            return extract_text_from_doc(file_path)
        
        elif file_extension == '.pptx':
            return extract_text_from_pptx(file_path)
        
        elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']:
            return extract_text_from_image(file_path)
        
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    except ImportError as e:
        # Re-raise ImportError to be handled by the caller
        raise e
    except Exception as e:
        raise Exception(f"Failed to extract text from {file_extension} file: {str(e)}")

def clean_extracted_text(text: str) -> str:
    """Clean and normalize extracted text, especially for Urdu content"""
    if not text:
        return ""
    
    # Decode HTML entities
    text = html.unescape(text)
    
    # Remove common OCR artifacts and noise
    # Remove standalone single characters that are likely OCR errors
    text = re.sub(r'\b[a-zA-Z]\b(?!\s+[اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے])', '', text)
    
    # Remove mixed language artifacts (English chars mixed with Urdu)
    # Keep only if they form meaningful words
    text = re.sub(r'\b[a-zA-Z]{1,2}\s+(?=[اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے])', '', text)
    
    # Remove standalone numbers that are OCR artifacts
    text = re.sub(r'\b\d{1,2}\b(?=\s+[اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے])', '', text)
    
    # Remove common OCR noise patterns
    noise_patterns = [
        r'\b(br|Fe|we|ei|Sule|بی گی)\b',  # Common OCR misreads
        r'\([^)]*\)',  # Remove parenthetical content that's often noise
        r'[{}\[\]<>]',  # Remove brackets and braces
        r'\s+[-_=+*#@$%^&]+\s+',  # Remove lines of special characters
        r'\b[a-zA-Z]+(\d+[a-zA-Z]*)+\b',  # Remove alphanumeric codes
    ]
    
    for pattern in noise_patterns:
        text = re.sub(pattern, ' ', text, flags=re.IGNORECASE)
    
    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single space
    text = text.strip()
    
    # If text is mostly Urdu, remove remaining English artifacts
    urdu_chars = len(re.findall(r'[اآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے]', text))
    total_chars = len(re.findall(r'[a-zA-Zاآبپتٹثجچحخدڈذرڑزژسشصضطظعغفقکگلمنںوہھیے]', text))
    
    if total_chars > 0 and urdu_chars / total_chars > 0.6:  # If 60% Urdu
        # Remove remaining English words that are likely OCR errors
        text = re.sub(r'\b[a-zA-Z]{1,3}\b', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
    
    return text

def get_supported_extensions():
    """Return list of supported file extensions based on available packages"""
    base_extensions = {'txt'}  # Always supported
    
    # Check for optional packages and add extensions if available
    try:
        import PyPDF2
        base_extensions.add('pdf')
    except ImportError:
        try:
            import pdfplumber
            base_extensions.add('pdf')
        except ImportError:
            pass
    
    try:
        from docx import Document
        base_extensions.add('docx')
    except ImportError:
        pass
    
    try:
        import textract
        base_extensions.add('doc')
    except ImportError:
        try:
            import docx2txt
            base_extensions.add('doc')
        except ImportError:
            pass
    
    try:
        from pptx import Presentation
        base_extensions.add('pptx')
    except ImportError:
        pass
    
    try:
        import pytesseract
        from PIL import Image
        base_extensions.update(['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'webp'])
    except ImportError:
        pass
    
    return base_extensions